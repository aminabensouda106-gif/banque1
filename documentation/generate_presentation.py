#!/usr/bin/env python3
"""Présentation PFA AmanaBank — gabarit type Support_Soutenance_PFE (EMSI / AmanaBank)."""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
IMG = ROOT / "overleaf" / "images"
UML = ROOT / "uml"
UML_UC = UML / "cas-utilisation"
UML_DATA = ROOT / "modele-donnees"
PRESENTATION_DIR = ROOT / "presentation"
OUT_PROJECT = PRESENTATION_DIR / "presentation-PFA-AmanaBank.pptx"
OUT_DOWNLOADS = Path(r"C:\Users\HP\Downloads\PFA.pptx")

ASSET_LOGOS = Path(
    r"C:\Users\HP\.cursor\projects\c-Users-HP-OneDrive-Ecole-Marocaine-des-Sciences-de-l-Ing-nieur-Bureau-Banque\assets"
)

DIAGRAM_SOURCES = [IMG, UML_UC, UML, UML_DATA]

# ── Palette (charte EMSI + AmanaBank, structure type soutenance PFE) ──
PRIMARY = RGBColor(0x00, 0x66, 0x33)       # EMSI green
PRIMARY_DARK = RGBColor(0x00, 0x4D, 0x27)
PRIMARY_LIGHT = RGBColor(0x4C, 0xAF, 0x7A)
ACCENT = RGBColor(0xA6, 0x89, 0x3A)        # AmanaBank gold
ACCENT_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
BROWN = RGBColor(0x4B, 0x36, 0x21)
BROWN_MID = RGBColor(0x6B, 0x53, 0x44)
BG_LIGHT = RGBColor(0xF9, 0xF6, 0xF1)
BG_PANEL = RGBColor(0xED, 0xE6, 0xDA)
BG_NAV_INACTIVE = RGBColor(0xE8, 0xF0, 0xEA)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x5C, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
LINE = RGBColor(0xC5, 0xD4, 0xC8)

STUDENTS = "BENSOUDA Amina & AFILAL Soraya"
ENCADRANTE = "Dr. ORCHI HOUDA"
FOOTER_LEFT = f"Soutenance PFA  ·  {STUDENTS}"
FOOTER_RIGHT = "EMSI Rabat  ·  AmanaBank"
SCHOOL_YEAR = "Année universitaire 2025 – 2026"
DEFENSE_DATE = "Soutenu le 30 Juin 2026"

NAV_SECTIONS = [
    "1. INTRODUCTION",
    "2. CONCEPTION",
    "3. RÉALISATION",
    "4. VALIDATION",
    "5. CONCLUSION",
]

PLAN_SECTIONS = [
    ("01", "Contexte du projet"),
    ("02", "Problématique"),
    ("03", "Objectifs"),
    ("04", "Méthodologie"),
    ("05", "Conception"),
    ("06", "Réalisation / démonstration"),
    ("07", "Résultats obtenus"),
    ("08", "Difficultés rencontrées"),
    ("09", "Conclusion et perspectives"),
]

SLIDE_COUNTER = {"n": 0, "total": 0}
PAGE_REFS: list = []


def refresh_diagram_svgs(*stems: str):
    """Regénère les SVG PlantUML et les copie vers overleaf/images."""
    script = ROOT.parent / "scripts" / "generate-plantuml-svgs.py"
    if script.exists():
        subprocess.run([sys.executable, str(script)], cwd=ROOT.parent, check=False)
    for stem in stems:
        for folder in (UML, UML_UC, UML_DATA):
            svg = folder / f"{stem}.svg"
            if svg.exists():
                shutil.copy2(svg, IMG / svg.name)
                break


def sync_diagram_pngs():
    """Copie les diagrammes PNG fournis (fond blanc) vers overleaf/images."""
    mapping = {
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_diageClass-b22adeaa-b90b-42b2-a4e0-4577f4073ef1.png": "diagramme-classes.png",
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_MCD-ab80d30a-280e-469f-8a68-6a85425d18d1.png": "MCD.png",
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_MLD-57148b67-bede-4f2e-829a-6083bfb395e8.png": "MLD.png",
    }
    IMG.mkdir(parents=True, exist_ok=True)
    for src_name, dst_name in mapping.items():
        src = ASSET_LOGOS / src_name
        if src.exists():
            shutil.copy2(src, IMG / dst_name)


def sync_logos():
    """Copie les logos fournis par l'étudiante vers overleaf/images."""
    mapping = {
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_EMSI-Emploi-Recrutement-750x375-42962c8e-e54d-47e3-8ea9-584ae137fce3.png": "emsi-logo.png",
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_ChatGPT_Image_18_juin_2026__03_07_20__1_-c0126a33-4375-4d29-ba93-762b4a84bef4.png": "amana-logo.png",
        "c__Users_HP_AppData_Roaming_Cursor_User_workspaceStorage_be4d552e63596ce272199eac8a557e84_images_YKIbMu6IAmr4d8dd3Afa-77c3ac9f-11dd-447e-bf45-8762d79cfeeb.png": "honoris-logo.png",
    }
    IMG.mkdir(parents=True, exist_ok=True)
    for src_name, dst_name in mapping.items():
        src = ASSET_LOGOS / src_name
        if src.exists():
            shutil.copy2(src, IMG / dst_name)


def pdf_to_png(pdf_path: Path, dpi: int = 300) -> Path | None:
    try:
        import fitz
    except ImportError:
        return None
    if not pdf_path.exists():
        return None
    out = PRESENTATION_DIR / "_tmp" / (pdf_path.stem + ".png")
    out.parent.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72), alpha=False)
    pix.save(str(out))
    doc.close()
    return _crop_whitespace(out)


def svg_to_png(svg_path: Path, target_width_px: int = 2800) -> Path | None:
    try:
        import fitz
    except ImportError:
        return None
    if not svg_path.exists():
        return None
    out = PRESENTATION_DIR / "_tmp" / (svg_path.stem + ".png")
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        doc = fitz.open(str(svg_path))
        page = doc[0]
        scale = target_width_px / max(page.rect.width, 1)
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        pix.save(str(out))
        doc.close()
    except Exception:
        return None
    return _crop_whitespace(out)


RASTER_EXT = (".png", ".jpg", ".jpeg")
VECTOR_EXT = (".svg", ".pdf")


def resolve_image(*candidates: str) -> Path | None:
    raster = set(RASTER_EXT)
    for name in candidates:
        base = name.replace(".pdf", "").replace(".svg", "").replace(".png", "")
        for folder in DIAGRAM_SOURCES:
            for ext in RASTER_EXT + VECTOR_EXT:
                p = folder / f"{base}{ext}"
                if not p.exists():
                    continue
                if p.suffix.lower() in raster:
                    return p
                if p.suffix.lower() == ".svg":
                    png = svg_to_png(p)
                    if png:
                        return png
                if p.suffix.lower() == ".pdf":
                    png = pdf_to_png(p)
                    if png:
                        return png
    return None


def _crop_whitespace(png_path: Path, padding: int = 14) -> Path | None:
    try:
        from PIL import Image, ImageChops
    except ImportError:
        return png_path
    img = Image.open(png_path).convert("RGB")
    bg_color = img.getpixel((0, 0))
    bg = Image.new("RGB", img.size, bg_color)
    diff = ImageChops.difference(img, bg)
    diff = diff.point(lambda p: 255 if p > 12 else 0)
    bbox = diff.getbbox()
    if not bbox:
        return png_path
    left, top, right, bottom = bbox
    left = max(0, left - padding)
    top = max(0, top - padding)
    right = min(img.width, right + padding)
    bottom = min(img.height, bottom + padding)
    img.crop((left, top, right, bottom)).save(png_path)
    return png_path


def image_aspect(path: Path) -> float:
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    return w / h if h else 1.0


def fit_in_box(aspect: float, max_w: float, max_h: float) -> tuple[float, float]:
    w, h = max_w, max_w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    return w, h


def add_picture_fitted(slide, path: Path, box_left: float, box_top: float, box_w: float, box_h: float):
    aspect = image_aspect(path)
    w, h = fit_in_box(aspect, box_w, box_h)
    left = box_left + (box_w - w) / 2
    top = box_top + (box_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w), height=Inches(h))
    return w, h


def add_image_frame(slide, left: float, top: float, width: float, height: float):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1.25)
    return frame


def nav_for_section(section_num: str) -> int:
    n = int(section_num)
    if n <= 3:
        return 0
    if n <= 5:
        return 1
    if n == 6:
        return 2
    if n <= 8:
        return 3
    return 4


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs: Presentation):
    SLIDE_COUNTER["n"] += 1
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def set_para(p, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, space_after=6):
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_after = Pt(space_after)


def add_rect(slide, left, top, width, height, fill, line=None, line_w=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s


def add_top_bar(slide):
    add_rect(slide, 0, 0, 13.333, 0.06, PRIMARY)


def add_left_bar(slide):
    add_rect(slide, 0, 0, 0.08, 7.5, PRIMARY)


def add_nav_tabs(slide, active_idx: int):
    add_rect(slide, 0, 0, 13.333, 0.288, BG_NAV_INACTIVE)
    tab_w = 13.333 / 5
    for i, label in enumerate(NAV_SECTIONS):
        x = i * tab_w
        active = i == active_idx
        bg = PRIMARY if active else BG_NAV_INACTIVE
        fg = WHITE if active else TEXT_MUTED
        tab = add_rect(slide, x, 0, tab_w, 0.4, bg)
        tf = tab.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        p = tf.paragraphs[0]
        set_para(p, label, 11 if active else 10, fg, active, PP_ALIGN.CENTER, 0)


def add_footer(slide, page_num: int | None = None, with_nav: bool = False):
    add_rect(slide, 0.65, 7.102, 12.033, 0.01, LINE)

    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.12), Inches(9.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_LEFT
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED

    rb = slide.shapes.add_textbox(Inches(7.5), Inches(7.12), Inches(4.5), Inches(0.3))
    rp = rb.text_frame.paragraphs[0]
    rp.text = FOOTER_RIGHT
    rp.font.size = Pt(9)
    rp.font.color.rgb = PRIMARY
    rp.alignment = PP_ALIGN.RIGHT

    emsi = resolve_image("emsi-logo.png")
    if emsi:
        slide.shapes.add_picture(str(emsi), Inches(9.93), Inches(7.05), height=Inches(0.27))
    amana = resolve_image("amana-logo.png")
    if amana:
        slide.shapes.add_picture(str(amana), Inches(11.85), Inches(7.05), height=Inches(0.22))

    if page_num is not None:
        pg = slide.shapes.add_textbox(Inches(10.35), Inches(7.05), Inches(2.3), Inches(0.35))
        p = pg.text_frame.paragraphs[0]
        p.text = str(page_num)
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT
        PAGE_REFS.append((page_num, p))


def add_slide_chrome(slide, active_nav: int | None = None, page_num: int | None = None):
    fill_bg(slide, WHITE)
    add_left_bar(slide)
    if active_nav is not None:
        add_nav_tabs(slide, active_nav)
    add_footer(slide, page_num, active_nav is not None)


def add_title_block(slide, title: str, rule_width: float = 5.1):
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.7 if True else 0.63), Inches(12.033), Inches(0.6))
    set_para(tb.text_frame.paragraphs[0], title.upper(), 24, PRIMARY, True)
    add_rect(slide, 0.65, 1.32, rule_width, 0.018, ACCENT)
    add_rect(slide, 0.65, 1.46, 12.033, 0.018, LINE)


def add_bullets(slide, items, left, top, width, height, size=16, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        level = 1 if item.startswith("  ") else 0
        p.text = item.strip()
        p.level = level
        p.font.size = Pt(size - (2 if level else 0))
        p.font.color.rgb = color
        p.font.bold = level == 0 and item.endswith(":")
        p.space_after = Pt(4)
    return box


def add_jury_table(slide):
    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(4.52), Inches(5.35), Inches(8.65), Inches(1.47))
    table = table_shape.table

    jury = [
        ("Prof. …", "Président du jury"),
        (ENCADRANTE, "Encadrante"),
        ("Prof. …", "Examinateur"),
    ]
    for r, (name, role) in enumerate(jury):
        c0 = table.cell(r, 0)
        c1 = table.cell(r, 1)
        c0.text = name
        c1.text = role
        for cell in (c0, c1):
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT


def slide_cover(prs, thanks: bool = False):
    slide = blank(prs)
    fill_bg(slide, BG_LIGHT)
    add_top_bar(slide)

    emsi = resolve_image("emsi-logo.png")
    if emsi:
        slide.shapes.add_picture(str(emsi), Inches(0.33), Inches(0.1), width=Inches(2.5))
    amana = resolve_image("amana-logo.png")
    if amana:
        slide.shapes.add_picture(str(amana), Inches(9.8), Inches(0.25), width=Inches(3.2))
    honoris = resolve_image("honoris-logo.png")
    if honoris:
        slide.shapes.add_picture(str(honoris), Inches(10.0), Inches(1.35), height=Inches(0.75))

    header = slide.shapes.add_textbox(Inches(2.86), Inches(0.29), Inches(7.4), Inches(1.3))
    tf = header.text_frame
    for i, (txt, sz, col) in enumerate([
        ("Soutenance du Projet de Fin d'Année", 16, PRIMARY),
        ("Pour l'obtention du diplôme d'Ingénieur d'État", 14, TEXT),
        ("Spécialité : Ingénierie Informatique et Réseaux", 14, BROWN_MID),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_para(p, txt, sz, col, i == 0, PP_ALIGN.CENTER, 4)

    panel = add_rect(slide, 1.05, 1.92, 11.2, 2.03, WHITE, LINE, 1)
    panel.shadow.inherit = False
    add_rect(slide, 0.99, 1.89, 2.24, 0.06, ACCENT)
    add_rect(slide, 0.99, 1.89, 0.06, 0.81, ACCENT)

    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(2.3), Inches(10.5), Inches(1.35))
    ttf = title_box.text_frame
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate([
        "CONCEPTION ET RÉALISATION D'UN SYSTÈME",
        "DE GESTION D'AGENCE BANCAIRE",
        "APPLICATION WEB AMANABANK",
    ]):
        p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        set_para(p, line, 22 if i == 1 else 18, BROWN if i == 1 else BROWN_MID, i == 1, PP_ALIGN.CENTER, 2)

    add_rect(slide, 10.02, 3.95, 2.24, 0.06, ACCENT)
    add_rect(slide, 12.25, 3.19, 0.06, 0.81, ACCENT)

    add_rect(slide, 0.55, 4.81, 12.2, 0.015, LINE)

    date_box = slide.shapes.add_textbox(Inches(5.42), Inches(4.22), Inches(2.5), Inches(0.35))
    set_para(date_box.text_frame.paragraphs[0], DEFENSE_DATE, 13, PRIMARY, True, PP_ALIGN.CENTER)

    left = slide.shapes.add_textbox(Inches(0.75), Inches(4.94), Inches(5.2), Inches(0.7))
    ltf = left.text_frame
    set_para(ltf.paragraphs[0], "Réalisé par :", 12, TEXT_MUTED, True)
    set_para(ltf.add_paragraph(), STUDENTS.replace(" & ", "\n"), 14, TEXT, True)

    jury_lbl = slide.shapes.add_textbox(Inches(4.52), Inches(4.93), Inches(6.15), Inches(0.35))
    set_para(jury_lbl.text_frame.paragraphs[0], "Devant le jury :", 12, TEXT_MUTED, True)
    add_jury_table(slide)

    year = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    set_para(year.text_frame.paragraphs[0], SCHOOL_YEAR, 13, PRIMARY, True, PP_ALIGN.CENTER)

    if thanks:
        overlay = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.3))
        tf = overlay.text_frame
        set_para(tf.paragraphs[0], "Merci pour votre attention", 32, PRIMARY, True, PP_ALIGN.CENTER, 8)
        set_para(tf.add_paragraph(), "Questions ?", 24, BROWN, True, PP_ALIGN.CENTER, 0)

    return slide


def slide_plan(prs):
    slide = blank(prs)
    fill_bg(slide, WHITE)
    add_left_bar(slide)
    add_top_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.58), Inches(12.033), Inches(0.6))
    set_para(tb.text_frame.paragraphs[0], "PRÉSENTATION : PLAN", 28, PRIMARY, True)
    add_rect(slide, 0.65, 1.22, 3.8, 0.018, ACCENT)
    add_rect(slide, 0.65, 1.36, 12.033, 0.018, LINE)

    plan_text = """INTRODUCTION
Contexte du projet | Problématique | Objectifs

CONCEPTION
Méthodologie | Acteurs & UML | Modèle de données | Règles de gestion

RÉALISATION
Stack technique | Architecture | Interfaces | Parcours démo

VALIDATION
Résultats obtenus | Difficultés rencontrées

CONCLUSION
Bilan & perspectives | Questions"""

    body = slide.shapes.add_textbox(Inches(1.18), Inches(1.67), Inches(10.0), Inches(5.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, block in enumerate(plan_text.split("\n\n")):
        lines = block.split("\n")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_para(p, lines[0], 16, PRIMARY, True, PP_ALIGN.LEFT, 6)
        if len(lines) > 1:
            set_para(tf.add_paragraph(), lines[1], 13, TEXT_MUTED, False, PP_ALIGN.LEFT, 14)

    note = slide.shapes.add_textbox(Inches(0.65), Inches(6.55), Inches(12.0), Inches(0.35))
    set_para(note.text_frame.paragraphs[0], "15 minutes de présentation  ·  5 minutes de questions", 12, GRAY)

    add_footer(slide, SLIDE_COUNTER["n"])


def slide_section_divider(prs, section_num: str, title: str):
    active = nav_for_section(section_num)
    slide = blank(prs)
    add_slide_chrome(slide, active, SLIDE_COUNTER["n"])

    add_rect(slide, 0.65, 1.72, 12.033, 0.018, LINE)

    macro = NAV_SECTIONS[active].split(". ", 1)[1]
    tb = slide.shapes.add_textbox(Inches(0.53), Inches(2.9), Inches(12.28), Inches(1.8))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_para(tf.paragraphs[0], macro.upper(), 40, PRIMARY, True, PP_ALIGN.CENTER, 8)
    set_para(tf.add_paragraph(), f"{section_num} — {title}", 22, TEXT_MUTED, False, PP_ALIGN.CENTER, 0)


def slide_content(prs, section_num: str, title: str, paragraphs: list[str] | None = None, bullets: list[str] | None = None):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}")

    top = 1.75
    if paragraphs:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.0), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            set_para(p, para, 16, TEXT, False, PP_ALIGN.JUSTIFY, 12)
    if bullets:
        add_bullets(slide, bullets, 0.65, top, 12.0, 5.2, 16, TEXT)


def slide_two_col(prs, section_num: str, title: str, bullets: list[str], image_names: list[str]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.5)
    add_bullets(slide, bullets, 0.65, 1.75, 5.8, 5.0, 15, TEXT)

    img = resolve_image(*image_names)
    if img:
        frame_left, frame_top, frame_w, frame_h = 6.55, 1.72, 6.35, 5.35
        add_image_frame(slide, frame_left, frame_top, frame_w, frame_h)
        add_picture_fitted(slide, img, frame_left + 0.08, frame_top + 0.08, frame_w - 0.16, frame_h - 0.16)


def slide_diagram(
    prs,
    section_num: str,
    title: str,
    image_names: list[str],
    caption: str = "",
    subtitle: str = "",
    hi_res: bool = False,
):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])

    header = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(12.3), Inches(0.5))
    set_para(header.text_frame.paragraphs[0], f"{section_num}. {title}", 20, PRIMARY, True)
    add_rect(slide, 0.65, 1.05, 2.5, 0.018, ACCENT)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.12), Inches(12.0), Inches(0.35))
        set_para(sub.text_frame.paragraphs[0], subtitle, 12, TEXT_MUTED)

    frame_top = 1.45 if subtitle else 1.2
    frame_left, frame_w, frame_h = 0.45, 12.45, 5.75 if subtitle else 5.75
    add_image_frame(slide, frame_left, frame_top, frame_w, frame_h)

    img = None
    for name in image_names:
        img = resolve_image(name)
        if img:
            break

    if img:
        add_picture_fitted(slide, img, frame_left + 0.08, frame_top + 0.08, frame_w - 0.16, frame_h - 0.16)
    else:
        err = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(1))
        set_para(err.text_frame.paragraphs[0], f"Diagramme introuvable : {', '.join(image_names)}", 14, GRAY, False, PP_ALIGN.CENTER)

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.65), Inches(6.95), Inches(12.0), Inches(0.35))
        set_para(cap.text_frame.paragraphs[0], caption, 11, TEXT_MUTED, False, PP_ALIGN.CENTER)


def slide_visual_context(prs, section_num: str, title: str, highlights: list[str], image_names: list[str]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.5)

    intro = slide.shapes.add_textbox(Inches(0.65), Inches(1.75), Inches(5.9), Inches(1.1))
    set_para(
        intro.text_frame.paragraphs[0],
        "Digitalisation bancaire · Agence AmanaBank · PFA EMSI Rabat 2025/2026",
        14,
        BROWN_MID,
        True,
    )

    for i, item in enumerate(highlights):
        y = 2.95 + i * 1.05
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(5.9), Inches(0.85))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BG_LIGHT if i % 2 == 0 else WHITE
        badge.line.color.rgb = PRIMARY_LIGHT
        badge.line.width = Pt(1)
        add_rect(slide, 0.65, y, 0.12, 0.85, PRIMARY)
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.12), Inches(5.5), Inches(0.6))
        set_para(tb.text_frame.paragraphs[0], item, 13, TEXT, False)

    img = resolve_image(*image_names)
    if img:
        add_image_frame(slide, 6.55, 1.72, 6.35, 5.35)
        add_picture_fitted(slide, img, 6.63, 1.8, 6.19, 5.19)


def slide_split_panels(
    prs,
    section_num: str,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
    question: str = "",
):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.5)

    panels = [
        (0.65, left_title, left_items, RGBColor(0xFD, 0xF2, 0xF2), RGBColor(0xC6, 0x28, 0x28)),
        (6.75, right_title, right_items, RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0xF9, 0xA8, 0x25)),
    ]
    for x, ptitle, items, bg, accent in panels:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(5.95), Inches(4.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = accent
        panel.line.width = Pt(1.5)
        add_rect(slide, x, 1.75, 5.95, 0.55, accent)
        th = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.82), Inches(5.65), Inches(0.45))
        set_para(th.text_frame.paragraphs[0], ptitle, 15, WHITE, True, PP_ALIGN.CENTER)
        add_bullets(slide, items, x + 0.2, 2.45, 5.55, 3.5, 14, TEXT)

    if question:
        qbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.25), Inches(12.05), Inches(0.65))
        qbox.fill.solid()
        qbox.fill.fore_color.rgb = PRIMARY
        qbox.line.fill.background()
        qt = slide.shapes.add_textbox(Inches(0.85), Inches(6.38), Inches(11.65), Inches(0.4))
        set_para(qt.text_frame.paragraphs[0], question, 13, WHITE, True, PP_ALIGN.CENTER)


def slide_goal_cards(prs, section_num: str, title: str, goals: list[tuple[str, str]]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.0)

    cols, gap = 2, 0.25
    cw = (12.05 - gap) / cols
    for i, (num, label) in enumerate(goals):
        col, row = i % cols, i // cols
        x = 0.65 + col * (cw + gap)
        y = 1.78 + row * 1.55
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1.25)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.28), Inches(0.75), Inches(0.75))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PRIMARY
        badge.line.fill.background()
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_para(btf.paragraphs[0], num, 14, WHITE, True, PP_ALIGN.CENTER, 0)

        lb = slide.shapes.add_textbox(Inches(x + 1.05), Inches(y + 0.2), Inches(cw - 1.2), Inches(0.95))
        lb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_para(lb.text_frame.paragraphs[0], label, 13, TEXT, False)


def slide_timeline(prs, section_num: str, title: str, phases: list[tuple[str, str]]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 5.5)

    add_rect(slide, 0.65, 3.35, 12.05, 0.06, ACCENT)
    n = len(phases)
    step_w = 12.05 / n
    for i, (phase, detail) in enumerate(phases):
        x = 0.65 + i * step_w + step_w / 2 - 0.35
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.15), Inches(0.7), Inches(0.7))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
        dot.line.fill.background()
        dtf = dot.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_para(dtf.paragraphs[0], str(i + 1), 16, WHITE, True, PP_ALIGN.CENTER, 0)

        ph = slide.shapes.add_textbox(Inches(0.65 + i * step_w), Inches(2.0), Inches(step_w - 0.1), Inches(0.9))
        ph.text_frame.word_wrap = True
        set_para(ph.text_frame.paragraphs[0], phase, 12, PRIMARY, True, PP_ALIGN.CENTER)

        det = slide.shapes.add_textbox(Inches(0.65 + i * step_w), Inches(4.0), Inches(step_w - 0.1), Inches(2.5))
        det.text_frame.word_wrap = True
        set_para(det.text_frame.paragraphs[0], detail, 11, TEXT_MUTED, False, PP_ALIGN.CENTER)

    note = slide.shapes.add_textbox(Inches(0.65), Inches(6.55), Inches(12.0), Inches(0.4))
    set_para(note.text_frame.paragraphs[0], "Chaque phase validée par mvn test avant de passer à la suivante.", 12, GRAY, False, PP_ALIGN.CENTER)


def slide_architecture_layers(prs, section_num: str, title: str):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 5.0)

    layers = [
        ("Navigateur", "Thymeleaf · Bootstrap 5 · Charte AmanaBank", RGBColor(0xE3, 0xF2, 0xFD), BROWN),
        ("Contrôleurs web", "MVC · /dashboard · /portal · /login", RGBColor(0xE8, 0xF5, 0xE9), PRIMARY),
        ("Services métier", "@Transactional · règles R1–R14 · audit", RGBColor(0xFF, 0xF8, 0xE1), ACCENT),
        ("Repositories JPA", "Spring Data · entités · requêtes", RGBColor(0xF3, 0xE5, 0xF5), BROWN_MID),
        ("PostgreSQL", "Flyway V1–V11 · banque_agence", RGBColor(0xEF, 0xEB, 0xE9), TEXT),
    ]
    x, w = 2.2, 8.9
    for i, (name, detail, bg, border) in enumerate(layers):
        y = 1.85 + i * 1.02
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.82))
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = border
        box.line.width = Pt(1.5)
        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.08), Inches(w - 0.4), Inches(0.65))
        tf = tb.text_frame
        set_para(tf.paragraphs[0], name, 14, border, True, PP_ALIGN.CENTER, 2)
        set_para(tf.add_paragraph(), detail, 11, TEXT_MUTED, False, PP_ALIGN.CENTER, 0)
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.45), Inches(y + 0.82), Inches(0.45), Inches(0.18))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    side = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.0), Inches(1.35), Inches(4.8))
    side.fill.solid()
    side.fill.fore_color.rgb = BG_PANEL
    side.line.color.rgb = PRIMARY
    st = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(1.15), Inches(4.4))
    st.text_frame.word_wrap = True
    set_para(st.text_frame.paragraphs[0], "Transversal", 12, PRIMARY, True, PP_ALIGN.CENTER, 8)
    set_para(st.text_frame.add_paragraph(), "Spring Security", 10, TEXT, False, PP_ALIGN.CENTER, 4)
    set_para(st.text_frame.add_paragraph(), "AuditService", 10, TEXT, False, PP_ALIGN.CENTER, 4)
    set_para(st.text_frame.add_paragraph(), "Notifications", 10, TEXT, False, PP_ALIGN.CENTER, 4)

    stats = slide.shapes.add_textbox(Inches(11.35), Inches(2.0), Inches(1.55), Inches(4.8))
    stf = stats.text_frame
    for line in ["9 entités", "80 tests", "20+ vues", "2 espaces"]:
        set_para(stf.paragraphs[0] if line == "9 entités" else stf.add_paragraph(), line, 11, PRIMARY, True, PP_ALIGN.LEFT, 10)


def slide_kpi_results(prs, section_num: str, title: str, kpis: list[tuple[str, str, str]], checks: list[str]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.5)

    n = len(kpis)
    kw = (12.05 - 0.2 * (n - 1)) / n
    for i, (value, label, color) in enumerate(kpis):
        x = 0.65 + i * (kw + 0.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.78), Inches(kw), Inches(1.55))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_rect(slide, x, 1.78, kw, 0.12, color)
        vb = slide.shapes.add_textbox(Inches(x), Inches(1.95), Inches(kw), Inches(0.7))
        set_para(vb.text_frame.paragraphs[0], value, 28, color, True, PP_ALIGN.CENTER, 0)
        lb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.65), Inches(kw - 0.2), Inches(0.55))
        set_para(lb.text_frame.paragraphs[0], label, 11, TEXT_MUTED, False, PP_ALIGN.CENTER)

    add_bullets(slide, checks, 0.65, 3.55, 12.0, 3.2, 15, TEXT)


def slide_issue_solution(prs, section_num: str, title: str, pairs: list[tuple[str, str]]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 6.5)

    for i, (issue, solution) in enumerate(pairs):
        y = 1.78 + i * 0.82
        h = 0.72
        prob = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(5.6), Inches(h))
        prob.fill.solid()
        prob.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xF2)
        prob.line.color.rgb = RGBColor(0xEF, 0x9A, 0x9A)
        pt = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(5.3), Inches(h - 0.15))
        set_para(pt.text_frame.paragraphs[0], f"⚠  {issue}", 11, TEXT, False)

        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(y + 0.18), Inches(0.55), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

        sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(y), Inches(5.65), Inches(h))
        sol.fill.solid()
        sol.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        sol.line.color.rgb = PRIMARY_LIGHT
        st = slide.shapes.add_textbox(Inches(7.2), Inches(y + 0.1), Inches(5.35), Inches(h - 0.15))
        set_para(st.text_frame.paragraphs[0], f"✓  {solution}", 11, TEXT, False)


def slide_rules_grid(prs, section_num: str, title: str, rules: list[tuple[str, str]]):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.5)

    cols, gap = 2, 0.22
    cw = (12.05 - gap) / cols
    for i, (code, desc) in enumerate(rules):
        col, row = i % cols, i // cols
        x = 0.65 + col * (cw + gap)
        y = 1.78 + row * 1.15
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHT if row % 2 == 0 else WHITE
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(y + 0.15), Inches(0.95), Inches(0.65))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = PRIMARY
        code_box.line.fill.background()
        ctf = code_box.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_para(ctf.paragraphs[0], code, 11, WHITE, True, PP_ALIGN.CENTER, 0)
        dt = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.15), Inches(cw - 1.35), Inches(0.65))
        dt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_para(dt.text_frame.paragraphs[0], desc, 12, TEXT, False)


def slide_stack_cards(
    prs,
    section_num: str,
    title: str,
    cards: list[tuple[str, list[str]]],
    header_colors: list[RGBColor] | None = None,
):
    slide = blank(prs)
    add_slide_chrome(slide, nav_for_section(section_num), SLIDE_COUNTER["n"])
    add_title_block(slide, f"{section_num}. {title}", 4.0)

    colors = header_colors or [PRIMARY, ACCENT, BROWN, PRIMARY_DARK]
    n = len(cards)
    cw = (12.0 - 0.2 * (n - 1)) / n
    for i, (card_title, items) in enumerate(cards):
        x = 0.65 + i * (cw + 0.2)
        hc = colors[i % len(colors)]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(cw), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = hc
        card.line.width = Pt(1.5)
        add_rect(slide, x, 1.75, cw, 0.55, hc)
        th = slide.shapes.add_textbox(Inches(x + 0.08), Inches(1.82), Inches(cw - 0.16), Inches(0.45))
        set_para(th.text_frame.paragraphs[0], card_title, 13, WHITE, True, PP_ALIGN.CENTER)
        add_bullets(slide, items, x + 0.12, 2.45, cw - 0.24, 4.1, 12, TEXT)


def build():
    sync_logos()
    sync_diagram_pngs()
    refresh_diagram_svgs("01-clients-comptes", "02-operations-supervision", "03-portail-client")
    SLIDE_COUNTER["n"] = 0
    PAGE_REFS.clear()

    tmp = PRESENTATION_DIR / "_tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    for old in tmp.glob("*.png"):
        try:
            old.unlink()
        except OSError:
            pass

    prs = new_prs()

    slide_cover(prs)
    slide_plan(prs)

    slide_section_divider(prs, "01", "Contexte du projet")
    slide_visual_context(
        prs, "01", "Contexte du projet",
        highlights=[
            "Digitalisation des services bancaires : enjeu majeur pour les agences marocaines.",
            f"Projet AmanaBank — PFA EMSI Rabat, encadré par {ENCADRANTE}.",
            "Simulation complète : accueil client, opérations, reporting et portail en ligne.",
        ],
        image_names=["landing.png"],
    )

    slide_section_divider(prs, "02", "Problématique")
    slide_split_panels(
        prs, "02", "Problématique",
        left_title="Constats terrain",
        left_items=[
            "Fiches clients papier ou fichiers dispersés",
            "Soldes suivis dans des tableurs",
            "Saisies manuelles sujettes aux erreurs",
            "Pas de journal centralisé des actions",
            "Aucun accès en ligne pour le client",
        ],
        right_title="Conséquences",
        right_items=[
            "Incohérences de solde et doubles saisies",
            "Traçabilité faible (agent non identifié)",
            "Supervision limitée pour le chef d'agence",
            "Délais de traitement allongés",
        ],
        question="Comment centraliser la gestion de l'agence tout en offrant un espace client sécurisé ?",
    )

    slide_section_divider(prs, "03", "Objectifs")
    slide_goal_cards(
        prs, "03", "Objectifs du projet",
        [
            ("1", "Analyser le besoin (cahier des charges PFA)"),
            ("2", "Modéliser le système (UML, MCD/MLD)"),
            ("3", "Architecture Spring Boot en couches"),
            ("4", "Noyau métier : clients, comptes, transactions"),
            ("5", "Extensions : factures, chéquier, notifications"),
            ("6", "Portail client + 80 tests d'intégration"),
        ],
    )

    slide_section_divider(prs, "04", "Méthodologie")
    slide_timeline(
        prs, "04", "Méthodologie de travail (2TUP)",
        [
            ("Conception", "UML · cas d'utilisation · MCD/MLD"),
            ("Bootstrap", "Spring Boot · Security · rôles"),
            ("Métier cœur", "Clients · comptes · transactions"),
            ("Reporting", "Dashboard KPI · audit · admin"),
            ("Extensions", "Factures · chéquier · notifications"),
            ("Canal client", "Portail en ligne · seed démo"),
        ],
    )

    slide_section_divider(prs, "05", "Conception")
    slide_stack_cards(
        prs, "05", "Acteurs du système",
        [
            ("Agent bancaire", ["Clients et comptes", "Opérations courantes", "Historique et chéquiers"]),
            ("Chef d'agence", ["Supervision et KPI", "Notifications", "Journal d'audit"]),
            ("Administrateur", ["Gestion du personnel", "Rôles et activation", "Supervision globale"]),
            ("Client bancaire", ["Portail en ligne", "Consultation comptes", "Notifications et chéquiers"]),
        ],
        header_colors=[PRIMARY, ACCENT, BROWN, PRIMARY_DARK],
    )
    slide_diagram(
        prs, "05", "Cas d'utilisation — clients et comptes",
        ["01-clients-comptes"],
        "Personnel agence · authentification · CRUD clients · gestion comptes (§7.1 – §7.2)",
        subtitle="Acteurs : administrateur, agent, chef d'agence",
    )
    slide_diagram(
        prs, "05", "Cas d'utilisation — opérations et supervision",
        ["02-operations-supervision"],
        "Transactions · factures · reporting · audit · alertes KPI (§7.3 – §7.5)",
        subtitle="Contrôle d'éligibilité : solde suffisant · compte actif",
    )
    slide_diagram(
        prs, "05", "Cas d'utilisation — portail client",
        ["03-portail-client"],
        "Espace /portal en lecture seule · CIN ou n° client · reçus PDF",
        subtitle="Acteur : client bancaire (rôle CLIENT)",
    )
    slide_diagram(
        prs, "05", "Diagramme de classes — domaine métier",
        ["diagramme-classes.png", "diagramme-classes"],
        "9 entités · hiérarchies User et Transaction · extensions métier",
    )
    slide_diagram(
        prs, "05", "Modèle conceptuel de données (MCD)",
        ["MCD.png", "MCD"],
        "Entités métier · cardinalités R1–R14 · packages par domaine",
    )
    slide_diagram(
        prs, "05", "Modèle logique de données (MLD)",
        ["MLD.png", "MLD"],
        "PostgreSQL banque_agence · 9 tables · Flyway V1–V11",
    )
    slide_rules_grid(
        prs, "05", "Règles de gestion clés",
        [
            ("R1–R2", "1 client → N comptes ; 1 compte → 1 client"),
            ("R3", "Retrait / paiement refusé si solde insuffisant"),
            ("R4", "Virement uniquement entre comptes actifs"),
            ("R5–R8", "Traçabilité · statuts compte · BCrypt · audit"),
            ("R9–R11", "Facture : référence obligatoire · reçu PDF"),
            ("R12–R14", "Chéquier : 1 PENDING max · sans impact solde"),
        ],
    )

    slide_section_divider(prs, "06", "Réalisation / démonstration")
    slide_stack_cards(
        prs, "06", "Stack technique",
        [
            ("Back-end", ["Java 21", "Spring Boot 3.4", "Spring Security + JPA", "PostgreSQL + Flyway", "OpenPDF"]),
            ("Front-end", ["Thymeleaf MVC", "Bootstrap 5.3", "Charte AmanaBank", "DM Sans, CSS custom"]),
            ("Qualité", ["JUnit 5 — 80 tests", "PlantUML", "Git + seed dev", "Maven"]),
        ],
        header_colors=[PRIMARY, ACCENT, BROWN],
    )
    slide_architecture_layers(prs, "06", "Architecture en couches")
    slide_two_col(
        prs, "06", "Interface — vitrine et connexion",
        [
            "Landing page institutionnelle (/)",
            "Charte marron · beige · or",
            "Connexion unique personnel + client",
            "Redirection : /dashboard ou /portal",
            "Identifiants démo sur l'écran",
        ],
        ["login.png", "landing.png"],
    )
    slide_two_col(
        prs, "06", "Interface — espace agence",
        [
            "Dashboard KPI (clients, comptes, volume)",
            "Alertes chéquier cliquables",
            "CRUD clients — recherche multicritère",
            "Opérations : dépôt, retrait, virement",
            "Paiement facture + reçu PDF",
            "Workflow chéquier et notifications",
        ],
        ["dashboard-agent.png"],
    )
    slide_two_col(
        prs, "06", "Interface — portail client",
        [
            "Consultation en lecture seule",
            "Synthèse patrimoine et comptes",
            "Historique + téléchargement PDF",
            "Notifications opérations agence",
            "Suivi commandes de chéquier",
            "Badge « Espace client »",
        ],
        ["portal-dashboard.png"],
    )
    slide_timeline(
        prs, "06", "Parcours de démonstration live (5–7 min)",
        [
            ("Landing", "Page d'accueil AmanaBank"),
            ("Agent", "Dashboard KPI · alertes"),
            ("Opération", "Dépôt · virement · facture"),
            ("Chéquier", "Commande · notification chef"),
            ("Client", "CIN MB654321 · portail"),
            ("Bilan", "Comptes · PDF · notifications"),
        ],
    )

    slide_section_divider(prs, "07", "Résultats obtenus")
    slide_kpi_results(
        prs, "07", "Résultats obtenus",
        [
            ("80", "Tests d'intégration", PRIMARY),
            ("11", "Migrations Flyway", ACCENT),
            ("9", "Entités JPA", BROWN),
            ("100%", "Cahier des charges", PRIMARY_DARK),
        ],
        checks=[
            "✓ Centralisation clients, comptes et transactions",
            "✓ Extensions livrées : factures, chéquier, notifications, portail",
            "✓ Règles métier R1–R14 appliquées et testées",
            "✓ Traçabilité complète (historique + AuditLog)",
            "✓ Sécurité BCrypt · rôles · portail lecture seule",
            "✓ Jeu de démo reproductible · 11 diagrammes UML · rapport complet",
        ],
    )

    slide_section_divider(prs, "08", "Difficultés rencontrées")
    slide_issue_solution(
        prs, "08", "Difficultés rencontrées et solutions",
        [
            ("PostgreSQL port 5432 occupé", "Basculement port 5433 + doc setup"),
            ("Cohérence transactionnelle des soldes", "@Transactional + @Version + tests"),
            ("Auth dual personnel / client sur /login", "UserDetailsServiceImpl composite"),
            ("Portail client sur base existante", "DemoPortalSync au démarrage (dev)"),
            ("Chéquier ≠ opération financière", "Entité CheckbookOrder sans impact solde"),
            ("Image bancaire professionnelle", "Charte AmanaBank complète (landing, KPI)"),
        ],
    )

    slide_section_divider(prs, "09", "Conclusion et perspectives")
    slide_split_panels(
        prs, "09", "Conclusion et perspectives",
        left_title="Bilan",
        left_items=[
            "Problématique résolue : gestion centralisée + portail client",
            "Cycle complet maîtrisé : analyse → UML → Spring Boot → tests",
            "Application déployable en démo (run-dev.bat :8081)",
        ],
        right_title="Perspectives",
        right_items=[
            "API REST documentée (Swagger)",
            "Export PDF relevés complets",
            "Alertes SMS/e-mail · authentification 2FA",
            "CI/CD Docker · conformité KYC",
        ],
        question="Merci pour votre attention — Questions ?",
    )

    slide_cover(prs, thanks=True)

    total = len(prs.slides)
    for num, para in PAGE_REFS:
        para.text = f"{num}/{total}"

    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    staging = OUT_PROJECT.with_suffix(".build.pptx")
    saved_project = staging
    for candidate in (staging, PRESENTATION_DIR / "presentation-PFA-AmanaBank-new.pptx"):
        try:
            prs.save(str(candidate))
            saved_project = candidate
            break
        except OSError as exc:
            if candidate == staging:
                print(f"WARN: save to {staging} failed ({exc}); trying alternate path", file=sys.stderr)
            else:
                raise
    try:
        if OUT_PROJECT.exists():
            OUT_PROJECT.unlink()
        if saved_project != OUT_PROJECT:
            saved_project.replace(OUT_PROJECT)
        saved_project = OUT_PROJECT
    except OSError:
        print(
            f"WARN: could not overwrite {OUT_PROJECT} (file may be open); kept {saved_project}",
            file=sys.stderr,
        )
    try:
        shutil.copy2(saved_project, OUT_DOWNLOADS)
        print(f"OK downloads: {OUT_DOWNLOADS}")
    except OSError as exc:
        print(f"WARN: copy to downloads failed: {exc}", file=sys.stderr)
    print(f"OK project : {saved_project}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    try:
        import fitz  # noqa: F401
    except ImportError:
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "pymupdf", "-q"])
    build()
